# ===== app/parsing/pptx_parser.py =====
import time
from pptx import Presentation
from app.parsing.base import BaseParser
from app.parsing.schema import ParsedDocument, PageInfo, Block

class PptxParser(BaseParser):
    name = "pptx_parser"
    version = "1.0"

    async def parse(self, document_id: str, content: bytes) -> ParsedDocument:
        import io
        t0 = time.perf_counter()
        prs = Presentation(io.BytesIO(content))
        blocks, pages, raw_text_parts = [], [], []
        order = 0

        for slide_idx, slide in enumerate(prs.slides, start=1):
            pages.append(PageInfo(page_number=slide_idx))
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    text = shape.text_frame.text
                    blocks.append(Block(block_id=f"{document_id}_s{slide_idx}_b{order}", type="paragraph",
                                          text=text, page=slide_idx, order=order))
                    raw_text_parts.append(text)
                    order += 1
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                note = slide.notes_slide.notes_text_frame.text
                blocks.append(Block(block_id=f"{document_id}_s{slide_idx}_notes", type="caption",
                                      text=note, page=slide_idx, order=order))
                order += 1

        return ParsedDocument(
            document_id=document_id, source_format="pptx", parser_name=self.name, parser_version=self.version,
            pages=pages, blocks=blocks, tables=[], images=[], raw_text="\n".join(raw_text_parts),
            reading_order_applied=True, ocr_used=False, processing_time_ms=(time.perf_counter() - t0) * 1000,
        )